#!/usr/bin/env python3
"""Presentation metadata and per-slide analysis. Always run this first to triage a PPT/PPTX file."""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Inches, Emu

# Import shared utilities from same directory
sys.path.insert(0, os.path.dirname(__file__))
from ppt_utils import ensure_pptx


def analyze(path: str) -> dict:
    pptx_path = ensure_pptx(path)
    prs = Presentation(pptx_path)

    slide_width_inches = round(prs.slide_width / 914400, 2) if prs.slide_width else 0
    slide_height_inches = round(prs.slide_height / 914400, 2) if prs.slide_height else 0
    aspect = "16:9" if abs((slide_width_inches / slide_height_inches) - (16/9)) < 0.1 else \
              "4:3" if abs((slide_width_inches / slide_height_inches) - (4/3)) < 0.1 else \
              "widescreen" if slide_width_inches / slide_height_inches > 1.5 else "standard"

    slides = []
    total_text_length = 0
    total_images = 0
    total_tables = 0
    total_charts = 0
    total_notes = 0
    empty_slides = 0

    for i, slide in enumerate(prs.slides):
        text_shapes = 0
        image_count = 0
        table_count = 0
        chart_count = 0
        other_count = 0
        slide_text_length = 0

        for shape in slide.shapes:
            shape_type = str(shape.shape_type)

            if shape.has_text_frame:
                text = shape.text_frame.text
                slide_text_length += len(text)
                if text.strip():
                    text_shapes += 1

            if shape.has_table:
                table_count += 1

            # Check for images/charts
            if 'PICTURE' in shape_type:
                image_count += 1
            elif 'CHART' in shape_type:
                chart_count += 1
            elif 'TABLE' in shape_type:
                table_count += 1
            elif 'GROUP' in shape_type:
                other_count += 1
            elif 'PLACEHOLDER' in shape_type:
                pass  # Counted under text shapes if it has text
            else:
                if not shape.has_text_frame and not shape.has_table:
                    other_count += 1

        # Check for speaker notes
        has_notes = False
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide
                notes_text = notes.notes_text_frame.text
                has_notes = bool(notes_text.strip())
        except Exception:
            has_notes = False

        if has_notes:
            total_notes += 1

        total_text_length += slide_text_length
        total_images += image_count
        total_tables += table_count
        total_charts += chart_count

        if slide_text_length < 20 and image_count == 0 and table_count == 0:
            empty_slides += 1

        # Try to get slide title
        title_text = ""
        try:
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()[:120]
        except Exception:
            pass

        # Get layout name
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        slides.append({
            "slide": i + 1,
            "layout": layout_name,
            "title": title_text,
            "text_shapes": text_shapes,
            "image_count": image_count,
            "table_count": table_count,
            "chart_count": chart_count,
            "other_shapes": other_count,
            "text_length": slide_text_length,
            "has_notes": has_notes,
            "notes_preview": notes_text[:100] if has_notes else "",
        })

    result = {
        "file": path,
        "format": os.path.splitext(path)[1].lower().lstrip('.'),
        "slide_count": len(slides),
        "slide_dimensions": f"{slide_width_inches}\" × {slide_height_inches}\" ({aspect})",
        "total_text_length": total_text_length,
        "total_images": total_images,
        "total_tables": total_tables,
        "total_charts": total_charts,
        "slides_with_notes": total_notes,
        "empty_slides": empty_slides,
        "slides": slides,
    }

    # Cleanup temp file if conversion happened
    if os.path.dirname(pptx_path).startswith("/tmp/pi-ppt-conv-"):
        import shutil
        shutil.rmtree(os.path.dirname(pptx_path), ignore_errors=True)

    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: ppt_info.py <path>", file=sys.stderr)
        sys.exit(1)

    result = analyze(sys.argv[1])
    print(json.dumps(result, indent=2))
