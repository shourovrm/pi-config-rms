#!/usr/bin/env python3
"""Extract text from presentation slides: titles, shapes, tables, and speaker notes."""

import sys
import os
import argparse
from pptx import Presentation

sys.path.insert(0, os.path.dirname(__file__))
from ppt_utils import ensure_pptx


def extract(path: str, slides_spec: str = "all", include_notes: bool = True,
            include_tables: bool = True, titles_only: bool = False) -> None:
    pptx_path = ensure_pptx(path)
    prs = Presentation(pptx_path)

    # Parse slide range
    if slides_spec.strip().lower() == "all":
        slide_indices = list(range(len(prs.slides)))
    else:
        slide_indices = set()
        total = len(prs.slides)
        for part in slides_spec.split(","):
            part = part.strip()
            if "-" in part:
                start, end = part.split("-", 1)
                start = max(1, int(start))
                end = min(total, int(end))
                slide_indices.update(range(start - 1, end))
            else:
                p = int(part) - 1
                if 0 <= p < total:
                    slide_indices.add(p)
        slide_indices = sorted(slide_indices)

    slide_list = list(prs.slides)

    for idx in slide_indices:
        slide = slide_list[idx]
        slide_num = idx + 1

        print(f"\n{'=' * 60}")
        print(f"=== Slide {slide_num} ===")

        # Layout
        try:
            layout = slide.slide_layout.name
            print(f"Layout: {layout}")
        except Exception:
            pass

        # Title
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass

        if title_shape and title_shape.text.strip():
            print(f"Title: {title_shape.text.strip()}")

        if titles_only:
            continue

        # Content shapes (skip the title if it was already shown)
        for shape in slide.shapes:
            if shape == title_shape:
                continue

            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shape_name = shape.name if shape.name else "Text Shape"
                    print(f"\n  [{shape_name}]")
                    for line in text.split('\n'):
                        print(f"    {line}")

            if include_tables and shape.has_table:
                table = shape.table
                if len(table.rows) > 0:
                    print(f"\n  [Table: {len(table.rows)}r × {len(table.columns)}c]")
                    for r_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace('\n', ' | ') for cell in row.cells]
                        print(f"    {'  |  '.join(cells)}")
                        if r_idx == 0 and len(table.rows) > 1:
                            print(f"    {'-' * 40}")

        # Speaker notes
        if include_notes:
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide
                    notes_text = notes.notes_text_frame.text.strip()
                    if notes_text:
                        print(f"\n  --- Speaker Notes ---")
                        print(f"  {notes_text}")
            except Exception:
                pass

    # Cleanup
    if os.path.dirname(pptx_path).startswith("/tmp/pi-ppt-conv-"):
        import shutil
        shutil.rmtree(os.path.dirname(pptx_path), ignore_errors=True)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text from presentation")
    parser.add_argument("path", help="Path to PPT/PPTX file")
    parser.add_argument("--slides", default="all",
                        help="Slide range: 'all', '1-5', '1,3,7', '3'")
    parser.add_argument("--no-notes", action="store_true",
                        help="Skip speaker notes")
    parser.add_argument("--no-tables", action="store_true",
                        help="Skip table extraction")
    parser.add_argument("--titles-only", action="store_true",
                        help="Extract only slide titles")
    args = parser.parse_args()

    extract(args.path, args.slides, not args.no_notes, not args.no_tables, args.titles_only)
