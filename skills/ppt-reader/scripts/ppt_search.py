#!/usr/bin/env python3
"""Search within presentation slides and speaker notes by keyword or regex."""

import sys
import os
import re
import argparse
from pptx import Presentation

sys.path.insert(0, os.path.dirname(__file__))
from ppt_utils import ensure_pptx


def search(path: str, pattern, context_lines: int = 2, literal: bool = False,
           max_matches: int = 50) -> int:
    pptx_path = ensure_pptx(path)
    prs = Presentation(pptx_path)

    total_matches = 0
    MAX_PREVIEW = 150

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Search in title
        try:
            title_shape = slide.shapes.title
            if title_shape:
                title_text = title_shape.text
                if pattern.search(title_text):
                    total_matches += 1
                    print(f"\n=== Slide {slide_num} [TITLE] ===")
                    print(f">>> {title_text[:MAX_PREVIEW]}")
        except Exception:
            pass

        # Search in all shapes
        for shape in slide.shapes:
            shape_name = shape.name or "Shape"

            if shape.has_text_frame:
                text = shape.text_frame.text
                lines = text.split('\n')

                for line_num, line in enumerate(lines):
                    if pattern.search(line):
                        total_matches += 1
                        print(f"\n=== Slide {slide_num}, [{shape_name}], line {line_num + 1} ===")

                        start = max(0, line_num - context_lines)
                        end = min(len(lines), line_num + context_lines + 1)
                        for j in range(start, end):
                            marker = ">>>" if j == line_num else "   "
                            print(f"{marker} {lines[j][:MAX_PREVIEW]}")

                        if total_matches >= max_matches:
                            print(f"\n--- Reached match limit ({max_matches}), stopping ---")
                            return total_matches

            if shape.has_table:
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        cell_text = cell.text
                        if pattern.search(cell_text):
                            total_matches += 1
                            print(f"\n=== Slide {slide_num}, [Table '{shape_name}'], row {r_idx + 1}, col {c_idx + 1} ===")
                            row_preview = [c.text[:80] for c in row.cells]
                            for ci, ct in enumerate(row_preview):
                                marker = ">>>" if ci == c_idx else "   "
                                print(f"{marker} [{ci + 1}] {ct}")

                            if total_matches >= max_matches:
                                print(f"\n--- Reached match limit ({max_matches}), stopping ---")
                                return total_matches

        # Search in speaker notes
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide
                notes_text = notes.notes_text_frame.text
                lines = notes_text.split('\n')

                for line_num, line in enumerate(lines):
                    if pattern.search(line):
                        total_matches += 1
                        print(f"\n=== Slide {slide_num} [SPEAKER NOTES], line {line_num + 1} ===")

                        start = max(0, line_num - context_lines)
                        end = min(len(lines), line_num + context_lines + 1)
                        for j in range(start, end):
                            marker = ">>>" if j == line_num else "   "
                            print(f"{marker} {lines[j][:MAX_PREVIEW]}")

                        if total_matches >= max_matches:
                            print(f"\n--- Reached match limit ({max_matches}), stopping ---")
                            return total_matches
        except Exception:
            pass

    # Cleanup
    if os.path.dirname(pptx_path).startswith("/tmp/pi-ppt-conv-"):
        import shutil
        shutil.rmtree(os.path.dirname(pptx_path), ignore_errors=True)

    return total_matches


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Search within presentation text")
    parser.add_argument("path", help="Path to PPT or PPTX file")
    parser.add_argument("query", help="Search query (regex or literal)")
    parser.add_argument("--context", type=int, default=2,
                        help="Context lines around match (default: 2)")
    parser.add_argument("--literal", action="store_true",
                        help="Treat query as literal string")
    parser.add_argument("--max-matches", type=int, default=50,
                        help="Maximum matches to report (default: 50)")
    args = parser.parse_args()

    if args.literal:
        pattern = re.compile(re.escape(args.query), re.IGNORECASE)
    else:
        try:
            pattern = re.compile(args.query, re.IGNORECASE)
        except re.error as e:
            print(f"Invalid regex: {e}", file=sys.stderr)
            sys.exit(1)

    ext = os.path.splitext(args.path)[1].lower()
    if ext not in ('.ppt', '.pptx'):
        print(f"ERROR: Unsupported format: {ext}. Use .ppt or .pptx", file=sys.stderr)
        sys.exit(1)

    total = search(args.path, pattern, args.context, args.literal, args.max_matches)

    if total == 0:
        print(f"No matches found for: {args.query}")
    else:
        print(f"\n--- {total} match(es) found ---")
